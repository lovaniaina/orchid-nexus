# /backend/create_presentation.py

import subprocess
import sys

# --- User-configurable content ---
YOUR_NAME = "[Your Name]"
YOUR_TITLE_CONTACT = "[Your Title / Contact Info]"
# --------------------------------

def install_and_import():
    """Checks for python-pptx and installs it if not found."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        print("'python-pptx' library is already installed.")
        return Presentation, Inches, Pt, RGBColor, PP_ALIGN
    except ImportError:
        print("Required library 'python-pptx' not found. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            print("Installation successful. Importing now.")
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            return Presentation, Inches, Pt, RGBColor, PP_ALIGN
        except Exception as e:
            print(f"Error: Failed to install 'python-pptx'. Please install it manually using 'pip install python-pptx'")
            print(f"Details: {e}")
            sys.exit(1)

def create_presentation():
    """Generates the 10-slide PowerPoint presentation."""
    Presentation, Inches, Pt, RGBColor, PP_ALIGN = install_and_import()
    
    prs = Presentation()

    # Slide 1: Title Slide
    slide1_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide1_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Orchid Nexus"
    subtitle.text = f"An Integrated Management Platform for NGOs\n\n{YOUR_NAME}\n{YOUR_TITLE_CONTACT}"

    # Slide 2: The Personal Connection
    slide2_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide2_layout)
    title = slide2.shapes.title
    title.text = "A Solution Born From Experience"
    body = slide2.shapes.placeholders[1].text_frame
    body.clear()
    p1 = body.add_paragraph()
    p1.text = "5+ Years Managing Humanitarian Projects in Madagascar"
    p1.level = 0
    p2 = body.add_paragraph()
    p2.text = "Bridging the Gap Between Field Operations and Strategic Goals"
    p2.level = 0

    # Slide 3: The Problem Statement
    slide3_layout = prs.slide_layouts[1]
    slide3 = prs.slides.add_slide(slide3_layout)
    title = slide3.shapes.title
    title.text = "The Crisis of Disconnected Data"
    body = slide3.shapes.placeholders[1].text_frame
    body.clear()
    p1 = body.add_paragraph()
    run1 = p1.add_run(); run1.text = "Program Data is a Black Hole: "; run1.font.bold = True
    run2 = p1.add_run(); run2.text = "Progress is invisible until it's too late."
    p2 = body.add_paragraph()
    run1 = p2.add_run(); run1.text = "Logistics are a Guessing Game: "; run1.font.bold = True
    run2 = p2.add_run(); run2.text = "Critical supplies are untracked, risking project failure."
    p3 = body.add_paragraph()
    run1 = p3.add_run(); run1.text = "Finances are Disconnected: "; run1.font.bold = True
    run2 = p3.add_run(); run2.text = "Budgets can't be reconciled against verified field work."

    # Slide 4: The Consequence
    slide4_layout = prs.slide_layouts[1]
    slide4 = prs.slides.add_slide(slide4_layout)
    title = slide4.shapes.title
    title.text = "The Result is Constant Crisis Management"
    body = slide4.shapes.placeholders[1].text_frame
    body.clear()
    p1 = body.add_paragraph(); p1.text = "Paper Forms"; p1.level = 0
    p2 = body.add_paragraph(); p2.text = "Spreadsheets"; p2.level = 0
    p3 = body.add_paragraph(); p3.text = "Email Chains"; p3.level = 0
    p4 = body.add_paragraph(); p4.text = "Lead to..."; p4.level = 1
    p5 = body.add_paragraph()
    run = p5.add_run(); run.text = "REACTIVE DECISIONS"; run.font.bold = True; run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    p5.level = 2

    # Slide 5: The Solution
    slide5_layout = prs.slide_layouts[1]
    slide5 = prs.slides.add_slide(slide5_layout)
    title = slide5.shapes.title
    title.text = "The Solution: A Single Source of Truth"
    body = slide5.shapes.placeholders[1].text_frame
    body.clear()
    p1 = body.add_paragraph(); p1.text = "Real-Time: Connects the field to the office, instantly."; p1.level = 0
    p2 = body.add_paragraph(); p2.text = "Role-Based: Provides a focused view for every team member."; p2.level = 0
    p3 = body.add_paragraph(); p3.text = "Integrated: Links program data, logistics, and deliverables."; p3.level = 0

    # Slide 6: Live Demo
    slide6_layout = prs.slide_layouts[6] # Blank layout
    slide6 = prs.slides.add_slide(slide6_layout)
    txBox = slide6.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Live Demonstration"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 7: How It Works
    slide7_layout = prs.slide_layouts[1]
    slide7 = prs.slides.add_slide(slide7_layout)
    title = slide7.shapes.title
    title.text = "The Technology Stack"
    body = slide7.shapes.placeholders[1].text_frame
    body.clear()
    p1 = body.add_paragraph(); run1 = p1.add_run(); run1.text = "Frontend: "; run1.font.bold = True; p1.add_run().text = "React (for a fast, interactive user experience)"
    p2 = body.add_paragraph(); run1 = p2.add_run(); run1.text = "Backend: "; run1.font.bold = True; p2.add_run().text = "Python & FastAPI (for robust, secure business logic)"
    p3 = body.add_paragraph(); run1 = p3.add_run(); run1.text = "Real-Time: "; run1.font.bold = True; p3.add_run().text = "WebSockets (for instant, two-way communication)"
    p4 = body.add_paragraph(); run1 = p4.add_run(); run1.text = "Security: "; run1.font.bold = True; p4.add_run().text = "JSON Web Tokens (for role-based data protection)"

    # Slide 8: The Roadmap
    slide8_layout = prs.slide_layouts[1]
    slide8 = prs.slides.add_slide(slide8_layout)
    title = slide8.shapes.title
    title.text = "What's Next"
    body = slide8.shapes.placeholders[1].text_frame
    body.clear()
    p1 = body.add_paragraph(); p1.text = "Phase 1: Program Monitoring & Field Data (✓ Complete)"; p1.level = 0
    p2 = body.add_paragraph(); p2.text = "Phase 2: Logistics Module (Inventory & Asset Tracking)"; p2.level = 0
    p3 = body.add_paragraph(); p3.text = "Phase 3: Finance Module (Budgeting & Expense Reporting)"; p3.level = 0
    
    # Slide 9: The Vision
    slide9_layout = prs.slide_layouts[1]
    slide9 = prs.slides.add_slide(slide9_layout)
    title = slide9.shapes.title
    title.text = "An Operating System for Impact"
    body = slide9.shapes.placeholders[1].text_frame
    body.clear()
    p1 = body.add_paragraph(); p1.text = "Integrate: Programs, Logistics, and Finance."; p1.level = 0
    p2 = body.add_paragraph(); p2.text = "Empower: Give managers the clarity and control to be proactive."; p2.level = 0
    p3 = body.add_paragraph(); p3.text = "Maximize: Enable NGOs to maximize their impact."; p3.level = 0

    # Slide 10: Thank You
    slide10_layout = prs.slide_layouts[5] # Title Only
    slide10 = prs.slides.add_slide(slide10_layout)
    title = slide10.shapes.title
    title.text = "Thank You"
    txBox = slide10.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{YOUR_NAME}\n{YOUR_TITLE_CONTACT}\n\nQuestions?"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)

    # Save the presentation
    file_name = "Orchid_Nexus_Presentation.pptx"
    prs.save(file_name)
    print(f"\nPresentation saved successfully as '{file_name}' in the current directory.")

if __name__ == "__main__":
    print("--- Orchid Nexus Presentation Generator ---")
    create_presentation()